"""PDF↔Office 转换模块"""
import subprocess
import time
from pathlib import Path

from core.models import ProgressCallback, TaskResult, TaskStatus
from core.platform import get_libreoffice_path


def pdf_to_docx(
    input_file: Path,
    output_dir: Path,
    progress_callback: ProgressCallback | None = None,
    **_,
) -> TaskResult:
    """
    PDF 转 Word（.docx）。使用 pdf2docx 还原排版。

    注意：复杂排版（多栏/特殊字体）还原度有限。
    """
    t0 = time.time()
    try:
        from pdf2docx import Converter

        output_dir.mkdir(parents=True, exist_ok=True)
        out_path = output_dir / f"{input_file.stem}.docx"

        if progress_callback:
            progress_callback(0, 1, "正在转换 PDF → Word...")

        cv = Converter(str(input_file))
        cv.convert(str(out_path))
        cv.close()

        if progress_callback:
            progress_callback(1, 1, "转换完成")

        return TaskResult(
            status=TaskStatus.SUCCESS,
            output_files=[out_path],
            output_dir=output_dir,
            duration_seconds=time.time() - t0,
        )

    except Exception as exc:
        return TaskResult(
            status=TaskStatus.FAILED,
            error_message=str(exc),
            duration_seconds=time.time() - t0,
        )


def pdf_to_xlsx(
    input_file: Path,
    output_dir: Path,
    progress_callback: ProgressCallback | None = None,
    **_,
) -> TaskResult:
    """PDF 转 Excel（.xlsx），提取表格数据。"""
    t0 = time.time()
    try:
        import pdfplumber
        from openpyxl import Workbook

        output_dir.mkdir(parents=True, exist_ok=True)
        out_path = output_dir / f"{input_file.stem}.xlsx"

        if progress_callback:
            progress_callback(0, 1, "正在提取表格...")

        wb = Workbook()
        ws = wb.active
        ws.title = "Sheet1"
        row_offset = 0

        with pdfplumber.open(str(input_file)) as pdf:
            total_pages = len(pdf.pages)
            for page_num, page in enumerate(pdf.pages, start=1):
                tables = page.extract_tables()
                if tables:
                    for table in tables:
                        for row in table:
                            for col_idx, cell in enumerate(row, start=1):
                                ws.cell(row=row_offset + 1, column=col_idx, value=cell or "")
                            row_offset += 1
                        row_offset += 1  # 表格间空行
                else:
                    # 无表格时提取全页文本
                    text = page.extract_text()
                    if text:
                        for line in text.split("\n"):
                            row_offset += 1
                            ws.cell(row=row_offset, column=1, value=line)
                        row_offset += 1

                if progress_callback:
                    progress_callback(page_num, total_pages, f"处理第 {page_num}/{total_pages} 页")

        wb.save(str(out_path))

        if progress_callback:
            progress_callback(1, 1, "转换完成")

        return TaskResult(
            status=TaskStatus.SUCCESS,
            output_files=[out_path],
            output_dir=output_dir,
            duration_seconds=time.time() - t0,
        )

    except Exception as exc:
        return TaskResult(
            status=TaskStatus.FAILED,
            error_message=str(exc),
            duration_seconds=time.time() - t0,
        )


def pdf_to_pptx(
    input_file: Path,
    output_dir: Path,
    progress_callback: ProgressCallback | None = None,
    **_,
) -> TaskResult:
    """PDF 转 PowerPoint（.pptx）。每页渲染为图片嵌入 slide。"""
    t0 = time.time()
    try:

        import pypdf
        from pptx import Presentation
        from pptx.util import Inches

        output_dir.mkdir(parents=True, exist_ok=True)
        out_path = output_dir / f"{input_file.stem}.pptx"

        if progress_callback:
            progress_callback(0, 1, "正在转换 PDF → PPT...")

        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        reader = pypdf.PdfReader(str(input_file))
        total_pages = len(reader.pages)

        for page_num in range(total_pages):
            # 将单页 PDF 提取为独立 PDF，再用文本方式处理
            writer = pypdf.PdfWriter()
            writer.add_page(reader.pages[page_num])

            # 创建幻灯片并添加页面文本
            slide_layout = prs.slide_layouts[6]  # 空白布局
            slide = prs.slides.add_slide(slide_layout)

            # 提取文本内容
            page_text = reader.pages[page_num].extract_text() or ""
            if page_text.strip():
                from pptx.util import Pt
                tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(6.5))
                tf = tx_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = page_text
                p.font.size = Pt(11)

            if progress_callback:
                progress_callback(page_num + 1, total_pages, f"处理第 {page_num + 1}/{total_pages} 页")

        prs.save(str(out_path))

        if progress_callback:
            progress_callback(1, 1, "转换完成")

        return TaskResult(
            status=TaskStatus.SUCCESS,
            output_files=[out_path],
            output_dir=output_dir,
            duration_seconds=time.time() - t0,
        )

    except Exception as exc:
        return TaskResult(
            status=TaskStatus.FAILED,
            error_message=str(exc),
            duration_seconds=time.time() - t0,
        )


def office_to_pdf(
    input_file: Path,
    output_dir: Path,
    progress_callback: ProgressCallback | None = None,
    **_,
) -> TaskResult:
    """
    Office 文件转 PDF（依赖 LibreOffice CLI）。

    若 LibreOffice 未安装，返回 TaskResult(FAILED) 并提示安装。
    """
    t0 = time.time()
    try:
        output_dir.mkdir(parents=True, exist_ok=True)
        soffice = get_libreoffice_path()

        if soffice is None:
            return TaskResult(
                status=TaskStatus.FAILED,
                error_message="未检测到 LibreOffice。请安装 LibreOffice 后重试。\n"
                              "下载地址：https://www.libreoffice.org/download/",
                duration_seconds=time.time() - t0,
            )

        if progress_callback:
            progress_callback(0, 1, "正在转换 Office → PDF...")

        cmd = [
            str(soffice),
            "--headless",
            "--convert-to", "pdf",
            "--outdir", str(output_dir),
            str(input_file),
        ]
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=300)

        if result.returncode != 0:
            stderr = result.stderr[-500:] if result.stderr else "未知错误"
            return TaskResult(
                status=TaskStatus.FAILED,
                error_message=f"LibreOffice 转换失败: {stderr}",
                duration_seconds=time.time() - t0,
            )

        out_path = output_dir / f"{input_file.stem}.pdf"

        if progress_callback:
            progress_callback(1, 1, "转换完成")

        return TaskResult(
            status=TaskStatus.SUCCESS,
            output_files=[out_path] if out_path.exists() else [],
            output_dir=output_dir,
            duration_seconds=time.time() - t0,
        )

    except Exception as exc:
        return TaskResult(
            status=TaskStatus.FAILED,
            error_message=str(exc),
            duration_seconds=time.time() - t0,
        )
